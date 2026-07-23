# -*- coding: utf-8 -*-
"""GOAI 初赛方案 PPT — 细化文案 + 平台实机截图。"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLATE = RGBColor(0x0F, 0x1C, 0x2E)
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_DARK = RGBColor(0x0A, 0x6B, 0x63)
INK = RGBColor(0x1A, 0x23, 0x32)
MUTED = RGBColor(0x5C, 0x6B, 0x7A)
LINE = RGBColor(0xD0, 0xD7, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)
CARD = RGBColor(0xE8, 0xF3, 0xF1)
AMBER = RGBColor(0xC4, 0x7A, 0x0A)

W, H = Inches(13.333), Inches(7.5)
SHOT = Path("docs/goai/out/screenshots")
TOTAL = 17


def _run(run, size=14, bold=False, color=INK, font="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(10), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = f"Testory · GOAI 新智基座 Agent Infra · 初赛方案 V0.3  |  {page}/{TOTAL}"
    _run(r, 10, color=MUTED)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), W, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_title(slide, text, top=0.28):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.5))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    _run(r, 24, bold=True, color=SLATE)


def add_subtitle(slide, text, top=0.78):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.38))
    r = box.text_frame.paragraphs[0].add_run()
    r.text = text
    _run(r, 12, color=TEAL_DARK)


def add_text(slide, left, top, width, height, lines, size=12, color=INK, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = line
        _run(r, size, bold=bold, color=color)


def add_card(slide, left, top, width, height, fill=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1)
    shape.adjustments[0] = 0.08
    return shape


def add_pill(slide, left, top, width, height, text, fill=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.adjustments[0] = 0.5
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _run(r, 11, bold=True, color=WHITE)


def style_table(table):
    for r, row in enumerate(table.rows):
        for cell in row.cells:
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(10 if r else 11)
                    run.font.bold = r == 0
                    run.font.color.rgb = WHITE if r == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = SLATE if r == 0 else (WHITE if r % 2 else RGBColor(0xEE, 0xF4, 0xF3))


def add_shot(slide, path, left, top, width, height):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        add_card(slide, left, top, width, height, CARD)
        add_text(slide, left + 0.2, top + height / 2 - 0.2, width - 0.4, 0.4, ["（截图缺失）"], size=12, color=MUTED)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ---- 01 Cover ----
    s = prs.slides.add_slide(blank)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.1), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = SLATE
    panel.line.fill.background()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), 0, Inches(0.12), H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    right = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.22), 0, Inches(8.12), H)
    right.fill.solid()
    right.fill.fore_color.rgb = SOFT
    right.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.3), Inches(4))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Testory"
    _run(r, 40, bold=True, color=WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "\n企业级多端联动质量保障\n多 Agent Infra"
    _run(r, 16, color=RGBColor(0xB8, 0xD4, 0xCF))
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "\n\nGOAI 新智基座｜Agent Infra\n赛题：复杂任务多 Agent 自主协同\n初赛方案 V0.3 · 含实机截图证据"
    _run(r, 11, color=TEAL)
    add_text(
        s,
        5.7,
        1.9,
        7.0,
        4.5,
        [
            "一条业务流，打通四端执行",
            "Web（本机 CDP）· Windows 桌面（UIA）",
            "Android（adb/插件）· REST API",
            "",
            "协同设计基点：AgentTeams",
            "能力沉淀：可复用 Skill 工程包",
            "执行运行时：Hermes（挂接各 Executor）",
            "开放定位：拟 Apache-2.0 开源",
            "",
            "队伍名称 / 成员（请填写）",
        ],
        size=14,
        color=SLATE,
    )

    # ---- 02 场景价值 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "真实场景：企业混合流回归，专家串测不可规模化")
    add_subtitle(s, "目标用户：测开 / QA / 业务质量团队；输入可为客诉、缺陷单或自然语言场景")
    add_card(s, 0.5, 1.35, 4.0, 5.2)
    add_text(
        s,
        0.7,
        1.55,
        3.6,
        4.8,
        [
            "痛点拆解",
            "",
            "· Web 下单失败常伴随桌面库存、",
            "  App 列表不同步，单端工具测不全",
            "· 变量（token/订单号）需人工抄录",
            "  到下一端，易漏易错",
            "· 验证码/扫码打断自动化后难续跑",
            "· 失败后只有截图碎片，缺证据等级",
            "  与可复用 Runbook",
            "· 现有 AI 方案多为单 Agent 对话，",
            "  难形成可审计 Production 闭环",
        ],
        size=12,
    )
    add_card(s, 4.7, 1.35, 4.0, 5.2, CARD)
    add_text(
        s,
        4.9,
        1.55,
        3.6,
        4.8,
        [
            "Testory 解题方式",
            "",
            "· Planner 把 NL/缺陷拆成 API→Web",
            "  →Desktop→Mobile 阶段计划",
            "· 各端 Executor 经 Skill 调真实工具：",
            "  CDP / UIA / adb / HTTP Runner",
            "· TestRunState 透传 vars 与证据索引",
            "· HitlGate 处理登录验证码后续跑",
            "· Verifier 输出强/弱/缺失证据报告",
            "· Skill 沉淀，避免一次性脚本",
        ],
        size=12,
    )
    add_card(s, 8.9, 1.35, 3.9, 5.2)
    add_text(
        s,
        9.1,
        1.55,
        3.5,
        4.8,
        [
            "可量化目标（Demo）",
            "",
            "· 跨端初诊/复现：",
            "  人工 30–60min → 目标 10–20min",
            "· 关键操作留痕：100%",
            "· L2 高风险动作审批：100%",
            "· 阶段变量透传：零人工抄录",
            "· 开源：平台核心 + Skills",
            "  可供企业/高校二次开发",
            "",
            "赛题对齐",
            "多 Agent · Skill · 工具契约",
            "证据 · 审批 · 开源成长",
        ],
        size=12,
    )
    add_footer(s, 2)

    # ---- 03 Demo 故事 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "主 Demo：客诉驱动的跨端复现与回归闭环")
    add_subtitle(s, "「10:15 后 Web 下单失败，桌面库存异常，App 订单列表未刷新」")
    add_card(s, 0.5, 1.3, 12.3, 1.5, CARD)
    add_text(
        s,
        0.7,
        1.45,
        11.9,
        1.2,
        [
            "触发：客诉/缺陷进入 → CapabilityProbe 探测本机 Edge/Chrome、桌面会话、adb 设备、API 可达性",
            "Planner：生成 CrossEndPlan（含 stages、vars_to_store、sync points、cleanup）并标注 L0/L1/L2",
            "执行：ApiHttp 查单 → WebBrowse 复现下单页 → WindowsDesktop 核库存 → AndroidMobile 刷列表 → Verifier 汇总",
        ],
        size=12,
    )
    steps = [
        ("T0", "客诉输入", "模糊描述\n+环境探针"),
        ("T1", "任务拆解", "Planner\n派单"),
        ("T2", "API 阶段", "查单/造数\n写 vars"),
        ("T3", "Web 阶段", "CDP 同屏\n可 HITL"),
        ("T4", "Desktop", "UIA/视觉\n核库存"),
        ("T5", "Mobile", "adb 动作\n列表断言"),
        ("T6", "验证沉淀", "证据评级\n复盘建议"),
    ]
    for i, (t, a, b) in enumerate(steps):
        left = 0.5 + i * 1.8
        add_card(s, left, 3.1, 1.7, 3.3)
        add_pill(s, left + 0.2, 3.3, 1.3, 0.35, t, fill=SLATE)
        add_text(s, left + 0.12, 3.85, 1.45, 2.3, [a, ""] + b.split("\n"), size=11, bold=True, color=SLATE)
    add_footer(s, 3)

    # ---- 04 截图 AI 自主测试 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实机证据①：AI 自主测试工作台（Web/桌面执行入口）")
    add_subtitle(
        s,
        "自然语言下达任务 → 本机启动 Edge/Chrome（CDP）或桌面直控 → 思考流程/动作/实时用例/报告四栏可观测",
    )
    add_shot(s, SHOT / "02_ai_test.png", 0.45, 1.25, 8.6, 5.4)
    add_card(s, 9.2, 1.25, 3.7, 5.4)
    add_text(
        s,
        9.4,
        1.45,
        3.3,
        5.0,
        [
            "界面能力对照赛题",
            "",
            "· 智能体服务：一键启动",
            "  Hermes/本机浏览器附着",
            "· 推理引擎可切换",
            "  （本地/云端模型）",
            "· 超时与「执行后生成用例」",
            "  把轨迹沉淀为可维护资产",
            "· 视觉识别开关：OCR/屏幕",
            "  观察兜底难定位控件",
            "· 右侧四栏 = 初赛可观测",
            "  雏形（决策·动作·用例·报告）",
            "",
            "对应 Agent：WebApiExecutor",
            "对应 Skill：WebBrowse",
            "对应闭环：工具调用+证据",
        ],
        size=11,
    )
    add_footer(s, 4)

    # ---- 05 截图 跨端 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实机证据②：跨端联动编排（Planner 原型）")
    add_subtitle(
        s,
        "NL → AI 分解为 API→Web→Mobile→Desktop 多阶段计划；模板含注册登录全链路、订单同步、数据一致性、移动/桌面对比",
    )
    add_shot(s, SHOT / "03_cross_end.png", 0.45, 1.25, 8.6, 5.4)
    add_card(s, 9.2, 1.25, 3.7, 5.4, CARD)
    add_text(
        s,
        9.4,
        1.45,
        3.3,
        5.0,
        [
            "为何这是关键证据",
            "",
            "· 已实现「任务拆解」产品面",
            "  = Planner Skill 的可演示形态",
            "· 明确四端阶段语义与",
            "  端到端变量透传、断言",
            "· 复赛将把此处分解结果",
            "  映射为 AgentTeams 任务图",
            "  与多角色派单事件",
            "",
            "诚实边界",
            "· 当前仍偏单回路编排",
            "· 尚未是对等多 Agent 对话",
            "· 初赛定位：协同 Spec 的",
            "  可运行前置能力",
        ],
        size=11,
    )
    add_footer(s, 5)

    # ---- 06 截图 AI中心 + 移动 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实机证据③：AI 中心模块化协作 + 移动端同步")
    add_subtitle(s, "用例设计 / 自主测试 / 自愈优化分模块；移动端支持 USB/无线配对、插件录制回放同步到 PC")
    add_shot(s, SHOT / "05_ai_hub.png", 0.4, 1.25, 6.3, 3.95)
    add_shot(s, SHOT / "04_mobile_testing.png", 6.85, 1.25, 6.0, 3.95)
    add_card(s, 0.4, 5.35, 12.5, 1.4, CARD)
    add_text(
        s,
        0.6,
        5.5,
        12.1,
        1.15,
        [
            "AI 中心：设计→执行→自愈 三段式，覆盖 Web / 接口 / 移动；跨端联动入口已产品化（对应赛题「经验沉淀 / 验证 / 工具接入」）。",
            "移动端同步：环境检测、设备连接、6 位配对码、助手插件；录制在手机完成、步骤回传 PC 编辑——对应 MobileExecutor + AndroidMobile Skill。",
        ],
        size=12,
    )
    add_footer(s, 6)

    # ---- 07 截图 登录 + 项目首页 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "实机证据④：本机独立模式 + 项目工作台")
    add_subtitle(s, "桌面壳登录强调数据本机保存；项目首页承载用例/执行入口，支撑企业向组织协作演进")
    add_shot(s, SHOT / "01_login.png", 0.4, 1.25, 5.9, 5.4)
    add_shot(s, SHOT / "06_projects.png", 6.5, 1.25, 4.0, 3.5)
    add_card(s, 6.5, 4.95, 6.4, 1.7, CARD)
    add_text(
        s,
        6.7,
        5.1,
        6.0,
        1.4,
        [
            "本机独立模式：数据保存在本机，降低企业数据外流顾虑，契合开源私有部署。",
            "项目工作台：用例与执行入口集中管理；审计/SSO 等为企业向同向能力（脚手架已有）。",
        ],
        size=11,
    )
    add_footer(s, 7)

    # ---- 08 架构 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "总体架构：AgentTeams 控制面 + Skill + 多端适配器")
    add_subtitle(s, "Agent 不裸调工具；Hermes 是端侧运行时，不是多 Agent 协同框架本身")
    layers = [
        ("① 入口层", "客诉 / 缺陷 / NL 场景 / CI webhook（复赛） / 平台 WebUI"),
        ("② AgentTeams 编排层（必须）", "角色拓扑 · 任务图 · 共享状态 · HITL/审批升级 · 状态机追踪"),
        ("③ 多 Agent 协同层", "Planner · WebApiExecutor · DesktopExecutor · MobileExecutor · Verifier"),
        ("④ Skill 能力层", "CrossEndDecompose · WebBrowse · ApiHttp · WindowsDesktop · AndroidMobile · RiskGuard…"),
        ("⑤ MCP / 适配器层", "CDP 附着 · Desktop Gateway · Mobile Gateway · HTTP Runner ·（可选）云 Skills"),
        ("⑥ 证据与治理层", "TestRunState · Trace/Log · 截图索引 · 审批记录 · Runbook/记忆（复赛）"),
    ]
    for i, (name, desc) in enumerate(layers):
        top = 1.25 + i * 0.85
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(12.3), Inches(0.75))
        bar.fill.solid()
        bar.fill.fore_color.rgb = SLATE if i == 1 else (CARD if i % 2 == 0 else WHITE)
        bar.line.color.rgb = LINE
        bar.adjustments[0] = 0.1
        color = WHITE if i == 1 else INK
        add_text(s, 0.7, top + 0.1, 4.0, 0.55, [name], size=13, bold=True, color=color)
        add_text(s, 4.8, top + 0.1, 7.7, 0.55, [desc], size=12, color=color)
    add_footer(s, 8)

    # ---- 09 Agents ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "五个职能 Agent：从拆解到验证的完整职责边界")
    add_subtitle(s, "赛题最低要求 ≥3；完整集 5 个。边界=职责分工（平台多端都能做，角色不越权兼岗）")
    agents = [
        ("Planner", "指挥", "NL/缺陷→CrossEndPlan\n派单·风险标注·重规划", "执行交由各端Executor"),
        ("WebApi", "执行", "CDP 浏览器+HTTP\n断言·HITL 续跑", "桌面/移动交由专责Agent"),
        ("Desktop", "执行", "UIA 优先·视觉降级\n窗口发现·变量衔接", "Web/API交由专责Agent"),
        ("Mobile", "执行", "adb/插件·dump/tap\n配对同步·列表断言", "Web/桌面交由专责Agent"),
        ("Verifier", "评判", "证据强/弱/缺失\nheal·复盘·沉淀建议", "证据不足不下武断根因"),
    ]
    for i, (n, tag, body, edge) in enumerate(agents):
        left = 0.4 + i * 2.55
        add_card(s, left, 1.3, 2.45, 5.3, CARD if i in (0, 4) else WHITE)
        add_pill(s, left + 0.25, 1.5, 1.95, 0.38, n, fill=SLATE if i == 0 else TEAL)
        add_text(s, left + 0.15, 2.1, 2.15, 0.35, [tag], size=11, bold=True, color=TEAL_DARK)
        add_text(s, left + 0.15, 2.55, 2.15, 2.4, body.split("\n"), size=11, color=SLATE)
        add_text(s, left + 0.15, 5.2, 2.15, 1.1, ["边界", edge], size=10, color=MUTED)
    add_footer(s, 9)

    # ---- 10 AgentTeams ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "AgentTeams 能力映射：初赛 Spec，复赛可加载运行")
    add_subtitle(s, "评审核验的是编排/拆解/上下文/协同/追踪如何落到框架，而非只提工具名")
    rows = [
        ("能力", "映射到本方案", "初赛交付", "复赛交付"),
        ("角色编排", "5 角色 Spec；最小 3 角色拓扑", "Identity 清单", "Spec 可加载"),
        ("任务拆解", "Planner + CrossEndDecompose→stages", "跨端页+API 证据", "任务图事件派单"),
        ("上下文传递", "TestRunState：vars/证据/审批令牌", "Schema 设计", "共享状态实现"),
        ("协同执行", "Executor+sync；Hermes=工具循环", "单端/跨端原型", "多角色并行/串行"),
        ("状态追踪", "test_run_id + trace_id 阶段机", "字段约定", "Trace 导出/看板"),
    ]
    table = s.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.35), Inches(12.3), Inches(4.5)).table
    for i, w in enumerate([1.8, 4.5, 3.0, 3.0]):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)
    add_text(
        s,
        0.5,
        6.05,
        12.3,
        0.7,
        ["红线表述：不把「单 Hermes 对话成功」写成「多 Agent 协同已完成」。AgentTeams=指挥部，Hermes=手。"],
        size=12,
        color=AMBER,
        bold=True,
    )
    add_footer(s, 10)

    # ---- 11 Skill ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Skill 工程体系：已有执行面 + 复赛治理/知识补齐")
    add_subtitle(s, "每个 Skill 具备输入输出、调用条件、失败处理、安全边界；源码位于 skills/bundled/")
    groups = [
        ("执行面（已开源资产）", [
            "WebBrowse — CDP 同屏+HITL",
            "ApiHttp — 造数/契约断言",
            "WindowsDesktop — UIA+视觉",
            "AndroidMobile — adb/插件",
            "CrossEndDecompose — NL→计划",
        ]),
        ("产品增强（已有）", [
            "UiDesignReview — UI 规范辅助",
            "AiExplore / AiDialog — 探索与对话式步骤",
            "CapabilityProbe — 开跑前能力矩阵",
        ]),
        ("治理类（复赛补齐）", [
            "HitlGate — 登录/验证码门禁",
            "RiskGuard — L0/L1/L2 白名单",
            "RecoveryVerify — 断言与恢复确认",
            "LocatorHeal — 定位失效候选",
        ]),
        ("知识类（复赛补齐）", [
            "RunbookRag — 规范/历史案例",
            "IncidentMemory — 决策记忆",
            "Postmortem — 复盘与沉淀清单",
        ]),
    ]
    for i, (title, items) in enumerate(groups):
        left = 0.4 + i * 3.2
        add_card(s, left, 1.3, 3.05, 5.3)
        add_text(s, left + 0.15, 1.5, 2.75, 0.55, [title], size=12, bold=True, color=TEAL_DARK)
        add_text(s, left + 0.15, 2.2, 2.75, 4.1, [f"· {x}" for x in items], size=11)
    add_footer(s, 11)

    # ---- 12 工具 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "工具契约：真实端连接层（MCP 或等价 Schema）")
    add_subtitle(s, "Skill 管能力抽象；适配器管协议/鉴权/错误/审计；Mock 与真机共用 Schema 以降迁移成本")
    tools = [
        ("Browser CDP Adapter", "launch_debug_browser → cdp_attach\n禁止默认 headless 偷跑\n验证码转 HitlGate"),
        ("Desktop Gateway", "端口化 UIA/视觉工具\ninprocess / gateway 模式\n前台窗口零配置发现"),
        ("Mobile Gateway", "adb forward · dump/tap/scroll\nscrcpy 预览协调\n助手插件配对码"),
        ("HTTP Runner", "method/url/headers/body\nassert.status · vars 透传\n密钥来自环境不入库"),
        ("AgentTeams（必须）", "角色·任务图·共享状态\n升级策略与追踪\n初赛 Spec / 复赛接入"),
        ("云 Skills（按需）", "云资源类能力可选用\n质量执行以自研 Skill 为主\n材料说明替换迁移成本"),
    ]
    for i, (n, d) in enumerate(tools):
        left = 0.4 + (i % 3) * 4.2
        top = 1.3 + (i // 3) * 2.7
        add_card(s, left, top, 4.0, 2.45)
        add_text(s, left + 0.2, top + 0.2, 3.6, 0.4, [n], size=13, bold=True, color=SLATE)
        add_text(s, left + 0.2, top + 0.75, 3.6, 1.5, d.split("\n"), size=12)
    add_footer(s, 12)

    # ---- 13 证据 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "上下文与证据分级：宁可报缺口，也不编造根因")
    add_subtitle(s, "共享 TestRunState；RAG/记忆复赛落地；满足赛题「记忆/共享状态/轨迹」中至少 2 项路径")
    cols = [
        ("输入归并", "客诉文本、API 响应、\n页面 DOM/截图、\n桌面控件树、移动 dump"),
        ("共享状态", "stages 状态机、vars、\n幂等键、审批令牌、\n证据索引"),
        ("检索（复赛）", "RunbookRag\nIncidentMemory\n失败模式召回"),
        ("证据评级", "强：时空一致可复现\n弱：部分采样不足\n缺：明确采集建议"),
        ("输出闭环", "pass/fail、heal 候选、\n复盘草稿、Skill 更新建议"),
    ]
    for i, (a, b) in enumerate(cols):
        left = 0.4 + i * 2.55
        add_card(s, left, 1.3, 2.45, 3.5, CARD if i >= 3 else WHITE)
        add_text(s, left + 0.15, 1.5, 2.15, 0.4, [a], size=13, bold=True, color=TEAL_DARK)
        add_text(s, left + 0.15, 2.1, 2.15, 2.5, b.split("\n"), size=11)
    add_card(s, 0.4, 5.0, 12.5, 1.7)
    add_text(
        s,
        0.6,
        5.2,
        12.1,
        1.4,
        [
            "示例：强证据 — Web /api/order/create 在 10:15 后持续 5xx，与客诉窗口一致且可复现；弱证据 — 桌面库存差异可见但缺变更审计；缺失 — 未采移动端网络 HAR → 输出采集建议而非「已定位根因」。",
            "初赛已具备：截图/动作日志/跨端计划结构。复赛补齐：统一 Trace 导出与向量检索。",
        ],
        size=12,
    )
    add_footer(s, 13)

    # ---- 14 安全 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "安全执行与可观测：L0–L2 + HITL + cleanup")
    add_subtitle(s, "企业级关键不是无限自动，而是可验证、可审批、可回滚、可审计")
    levels = [
        ("L0 只读", "探测、截图、日志拉取、能力探针", "默认可自动"),
        ("L1 低风险", "页面点击、查询类 API、常规断言", "白名单自动"),
        ("L2 高风险", "清数据、装/卸 APK、改配置、写生产", "审批令牌+人工"),
        ("HITL", "登录、滑块、扫码、验证码", "暂停→用户完成→续跑"),
        ("Cleanup", "跨端 cleanup stage 清理造数/恢复现场", "失败亦尽量执行"),
    ]
    for i, (n, d, p) in enumerate(levels):
        left = 0.4 + i * 2.55
        add_card(s, left, 1.3, 2.45, 3.3)
        add_pill(s, left + 0.2, 1.5, 2.05, 0.4, n, fill=SLATE if i < 3 else TEAL)
        add_text(s, left + 0.15, 2.15, 2.15, 2.2, [d, "", p], size=11)
    add_card(s, 0.4, 4.85, 12.5, 1.85, CARD)
    add_text(
        s,
        0.6,
        5.05,
        12.1,
        1.5,
        [
            "可观测：每次拆解 / Skill / 工具 / LLM 调用写入 Trace+Log，Metrics 含阶段时延、工具成功率、Token（若云模型）；关联 test_run_id。",
            "实现路径：复赛优先结构化 JSON Trace（可迁移 OTel GenAI / AgentLoop）；平台已有思考流程·执行动作·报告栏作为 UI 侧雏形。",
        ],
        size=12,
    )
    add_footer(s, 14)

    # ---- 15 路线 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "初赛边界 · 复赛闭环 · 与产品北极星同向")
    add_subtitle(s, "北极星：企业级多端（PC/Web/移动）联动自动化测试；比赛逼出多 Agent Infra 控制面")
    phases = [
        (
            "初赛 V0.3 · 可评审",
            [
                "完整方案 + Identity/Skill",
                "AgentTeams 映射表",
                "实机截图：AI 测试/跨端/移动/中心",
                "多端执行面与 HITL 已可用",
                "开源计划与诚实披露",
            ],
        ),
        (
            "复赛 V0.5 · 可运行",
            [
                "AgentTeams 接入最小 3 角色",
                "一条跨端故事端到端 Trace",
                "HITL + L2 审批各演示一次",
                "Desktop 或 Mobile 进主链路",
                "Skill Schema + MCP 样例包",
            ],
        ),
        (
            "长期 · 可成长",
            [
                "Apache-2.0 社区迭代",
                "Skill/MCP 生态贡献指南",
                "评测集与失败回放",
                "企业审计/权限加深",
                "开源新锐/影响力专项路径",
            ],
        ),
    ]
    for i, (t, items) in enumerate(phases):
        left = 0.4 + i * 4.2
        add_card(s, left, 1.3, 4.0, 5.3)
        add_text(s, left + 0.2, 1.5, 3.6, 0.55, [t], size=14, bold=True, color=TEAL_DARK)
        add_text(s, left + 0.2, 2.3, 3.6, 4.0, [f"· {x}" for x in items], size=13)
    add_footer(s, 15)

    # ---- 16 开源 ----
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "开源披露：非商业闭源，按可复现工程参赛")
    add_subtitle(s, "协议拟 Apache-2.0；核心仓 + Skills + Demo；密钥与未脱敏数据永不公开")
    add_card(s, 0.5, 1.3, 6.1, 5.3)
    add_text(
        s,
        0.7,
        1.5,
        5.7,
        5.0,
        [
            "开源范围",
            "· 平台核心（本地运行、AI 任务、跨端原型）",
            "· skills/bundled/* 全部执行面 Skill",
            "· 桌面/移动/Web 适配层与文档",
            "· docs/goai 方案与复赛 Demo 包",
            "",
            "依赖披露",
            "· Hermes：端侧工具运行时",
            "· Playwright / pywinauto / adb 等",
            "· 可选云 LLM（可替换 Ollama 本地）",
            "",
            "可复现",
            "· README · .env.example · 合成样例数据",
            "· 本 PPT 截图均来自本机 127.0.0.1:5000",
        ],
        size=12,
    )
    add_card(s, 6.9, 1.3, 5.9, 5.3, CARD)
    add_text(
        s,
        7.1,
        1.5,
        5.5,
        5.0,
        [
            "诚实披露（防红线）",
            "",
            "· 基于已有 Testory 工程继续演进",
            "· 新增重点：多 Agent 角色 + AgentTeams 映射",
            "· 现状代码主路径仍是单 Hermes 执行面",
            "  + 跨端阶段编排原型",
            "· 对等多 Agent 协作：初赛完成设计，",
            "  复赛交付可运行链路",
            "· 不刷 Star / 不伪造评测",
            "· 允许商业 API，但给出本地替代路径",
        ],
        size=12,
    )
    add_footer(s, 16)

    # ---- 17 收尾 ----
    s = prs.slides.add_slide(blank)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = SLATE
    panel.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.7), W, Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL_DARK
    bar.line.fill.background()
    box = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(12), Inches(3.2))
    tf = box.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = "Testory"
    _run(r, 36, bold=True, color=WHITE)
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = (
        "\n让复杂跨端质量保障，从专家手工串测，"
        "走向可协作、可治理、可复用的 Agent Infra。"
    )
    _run(r, 16, color=RGBColor(0xB8, 0xD4, 0xCF))
    p = tf.add_paragraph()
    r = p.add_run()
    r.text = "\nWeb · Windows · Android · API  —  同一 TestRun，共享状态与证据链"
    _run(r, 13, color=TEAL)
    add_text(
        s,
        0.7,
        5.95,
        12,
        1.3,
        [
            "01 多端联动实机　02 多 Agent 方案闭环　03 Skill 工程化　04 开源可审计",
            "下一步：填写队伍信息 → 提交初赛材料 → 按 docs/goai/semifinal_backlog.md 接入 AgentTeams",
        ],
        size=13,
        color=WHITE,
    )

    out_dir = Path("out")
    out_dir.mkdir(parents=True, exist_ok=True)
    goai_out = Path("docs/goai/out")
    goai_out.mkdir(parents=True, exist_ok=True)
    name = "Testory_GOAI_初赛方案.pptx"
    alt = "Testory_GOAI_初赛方案_v03.pptx"
    targets = [goai_out / name, goai_out / alt, out_dir / alt]
    # prefer overwriting primary out path when not locked
    targets.insert(0, out_dir / name)
    saved = []
    for path in targets:
        try:
            prs.save(str(path))
            saved.append(path.resolve())
        except PermissionError:
            print("locked, skip", path)
    for p in saved:
        print(p)
    if not saved:
        raise SystemExit("all save targets locked; close PPTX and retry")


if __name__ == "__main__":
    build()
